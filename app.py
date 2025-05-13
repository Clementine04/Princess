from flask import Flask, render_template, request, redirect, url_for, flash, send_file, jsonify, session, Response
import os
import threading
import json
import re  # For cleaning JSON strings
import google.generativeai as genai  # For AI-powered quiz and flashcard generation
from jinja2 import TemplateNotFound
from io import BytesIO
import pyttsx3
import random

# Import modules for file processing
from PyPDF2 import PdfReader
import docx
from pptx import Presentation

# Import database models
from models.database import db, Subject, Topic, Flashcard, QuizResult


API_KEY = os.environ.get("GOOGLE_GENAI_API_KEY", "AIzaSyDdSvIewEdADKXSWYriZJGqONj-KCXdsaU")
genai.configure(api_key=API_KEY)
MODEL_NAME = "gemini-2.0-flash"

app = Flask(__name__)
app.secret_key = 'your_secret_key'

# Configure the SQLAlchemy database
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///princess.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

# Initialize the database with the app
db.init_app(app)

# Create the database tables
with app.app_context():
    db.create_all()

questions = []
answers = []
options = []
question_types = []


def parse_ai_content(ai_generated_content, requested_count):
    """
    Parses AI output text into quiz data with improved answer handling and validation.
    Processes JSON-formatted quiz content and ensures proper structure and question types.
    """
    questions = []
    answers = []
    options = []
    question_types = []

    # Check if content is JSON formatted
    try:
        cleaned_text = clean_json_output(ai_generated_content)
        json_data = json.loads(cleaned_text)
        
        if isinstance(json_data, list):
            valid_questions = []
            
            for item in json_data:
                if "question" in item and "answer" in item:
                    q_type = item.get("type", "ID").strip().upper()
                    q_text = item["question"].strip()
                    answer = item["answer"].strip()
                    
                    # Skip empty questions or answers
                    if not q_text or not answer:
                        continue
                        
                    # Process based on question type
                    if q_type == "MC" and "options" in item and len(item["options"]) >= 2:
                        # Clean and validate multiple choice options
                        opts = [opt.strip() for opt in item["options"] if opt.strip()]
                        
                        # Ensure the correct answer is in the options
                        if answer not in opts:
                            # Add the answer to options if missing
                            opts.append(answer)
                            
                        # Ensure we have exactly 4 options
                        while len(opts) < 4:
                            opts.append(f"Option {len(opts) + 1}")
                            
                        # Limit to 4 options if more were provided
                        if len(opts) > 4:
                            # Keep the correct answer and sample 3 other options
                            other_options = [opt for opt in opts if opt != answer]
                            if len(other_options) > 3:
                                other_options = random.sample(other_options, 3)
                            opts = [answer] + other_options
                        
                        # Shuffle options to randomize position of correct answer
                        random.shuffle(opts)
                        
                        valid_questions.append({
                            "question": q_text,
                            "answer": answer,
                            "options": opts,
                            "type": "MC"
                        })
                        
                    elif q_type == "TF":
                        # Standardize True/False format
                        if answer.lower() in ["true", "t", "yes"]:
                            answer = "True"
                        elif answer.lower() in ["false", "f", "no"]:
                            answer = "False"
                        else:
                            # Skip invalid T/F answers
                            continue
                            
                        valid_questions.append({
                            "question": q_text,
                            "answer": answer,
                            "options": ["True", "False"],
                            "type": "TF"
                        })
                        
                    elif q_type == "ENUM":
                        # Clean enumeration answers
                        if "," in answer:
                            # Format and clean the comma-separated list
                            enum_items = [item.strip() for item in answer.split(",") if item.strip()]
                            if enum_items:
                                answer = ", ".join(enum_items)
                                valid_questions.append({
                                    "question": q_text,
                                    "answer": answer,
                                    "options": [],
                                    "type": "ENUM"
                                })
                        else:
                            # Single item doesn't need to be an enumeration
                            valid_questions.append({
                                "question": q_text,
                                "answer": answer,
                                "options": [],
                                "type": "ID"
                            })
                            
                    else:  # Default to ID type
                        # Validate and clean identification answers
                        answer = answer.strip()
                        if len(answer.split()) <= 5:  # Keep answers concise
                            valid_questions.append({
                                "question": q_text,
                                "answer": answer,
                                "options": [],
                                "type": "ID"
                            })
            
            # Truncate to requested count if needed
            if len(valid_questions) > requested_count:
                valid_questions = valid_questions[:requested_count]
                
            # Extract data from valid questions
            for q in valid_questions:
                questions.append(q["question"])
                answers.append(q["answer"])
                options.append(q["options"])
                question_types.append(q["type"])
            
            if questions:  # If we have valid questions, return them
                return {
                    'questions': questions,
                    'answers': answers,
                    'options': options,
                    'question_types': question_types
                }
    except (json.JSONDecodeError, ValueError) as e:
        print("Error parsing JSON:", e)
        # If not JSON, fall back to text parsing
        pass

    # Legacy text parsing method if JSON parsing failed
    lines = ai_generated_content.split("\n")
    current_question = None
    choice_lines = []

    def store_question(q, choices):
        if not q or not choices:
            return
            
        if len(choices) >= 4:  # Multiple choice
            correct = choices[0]
            opts = choices[:4]  # Take first 4 choices
            random.shuffle(opts)
            questions.append(q)
            answers.append(correct)
            options.append(opts)
            question_types.append("MC")
        else:
            if choices and choices[0].strip().lower() in ["true", "false", "true or false"]:
                questions.append(q)
                answers.append(choices[0].capitalize())
                options.append(["True", "False"])
                question_types.append("TF")
            else:
                answer = choices[0].strip() if choices else ""
                if not answer:
                    return
                    
                questions.append(q)
                answers.append(answer)
                options.append([])
                if "," in answer:
                    question_types.append("ENUM")
                else:
                    question_types.append("ID")

    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.lower().startswith("q:"):
            if current_question and choice_lines:
                store_question(current_question, choice_lines)
            current_question = line[2:].strip()
            choice_lines = []
        elif line.lower().startswith(("a:", "b:", "c:", "d:")) and current_question:
            choice_lines.append(line[2:].strip())
    if current_question and choice_lines:
        store_question(current_question, choice_lines)

    # Truncate to requested_count if needed
    if len(questions) > requested_count:
        questions = questions[:requested_count]
        answers = answers[:requested_count]
        options = options[:requested_count]
        question_types = question_types[:requested_count]

    return {
        'questions': questions,
        'answers': answers,
        'options': options,
        'question_types': question_types
    }


def build_quiz_prompt(module_text, exam_type, difficulty, num_q):
    """
    Builds an advanced prompt to instruct the AI to generate quiz questions
    with proper difficulty levels and structured outputs.
    """
    # Define difficulty-specific attributes
    difficulty_attributes = {
        "Easy": {
            "desc": "basic terminology and straightforward facts",
            "depth": "recall and simple understanding",
            "complexity": "direct questions with clear, unambiguous answers",
            "examples": {
                "MC": "What is the primary function of the heart?",
                "TF": "The heart pumps blood throughout the body.",
                "ID": "The main organ responsible for pumping blood.",
                "ENUM": "List the four chambers of the heart."
            }
        },
        "Medium": {
            "desc": "application of concepts and connecting ideas",
            "depth": "understanding relationships and moderate analysis",
            "complexity": "questions requiring some thought but with definite answers",
            "examples": {
                "MC": "Which of the following best describes the relationship between the heart and lungs?",
                "TF": "The right ventricle pumps oxygenated blood to the body tissues.",
                "ID": "The condition where heart valves don't close properly, causing blood to flow backward.",
                "ENUM": "List three factors that can increase the risk of heart disease."
            }
        },
        "Hard": {
            "desc": "deep analysis, synthesis of concepts, and critical thinking",
            "depth": "evaluation of complex scenarios and application in novel contexts",
            "complexity": "challenging questions that test comprehensive understanding",
            "examples": {
                "MC": "Which physiological mechanism would most effectively compensate for reduced cardiac output in a patient with heart failure?",
                "TF": "Frank-Starling's law explains why increased end-diastolic volume can initially improve cardiac output in heart failure.",
                "ID": "The compensatory mechanism where the heart enlarges to maintain stroke volume despite weakened contractions.",
                "ENUM": "List four pathophysiological changes that occur during cardiac remodeling after myocardial infarction."
            }
        }
    }

    # Get attributes for the selected difficulty
    diff_attr = difficulty_attributes[difficulty]
    
    base_prompt = f"""Generate {num_q} high-quality {exam_type} questions at {difficulty} difficulty level based on the content below.

DIFFICULTY LEVEL GUIDELINES - {difficulty.upper()} LEVEL:
• Focus on: {diff_attr["desc"]}
• Cognitive depth: {diff_attr["depth"]}
• Complexity: {diff_attr["complexity"]}

FORMAT REQUIREMENTS:
Return results in a JSON array with this exact structure:
[
  {{
    "question": "Clear question text ending with a question mark?",
    "answer": "Precise answer - specific terms only",
    "type": "QUESTION_TYPE",
    "options": ["Option A", "Option B", "Option C", "Option D"] // For multiple choice only
  }}
]

QUESTION TYPES:
"""

    # Add specific instructions based on exam type
    if exam_type == "Multiple Choice" or exam_type == "Mixed":
        base_prompt += f"""• MC (Multiple Choice):
  - One unambiguously correct answer among 4 options
  - Plausible distractors based on common misconceptions
  - {difficulty} example: "{diff_attr["examples"]["MC"]}"
  - Format: "options": ["Correct Answer", "Distractor 1", "Distractor 2", "Distractor 3"]

"""

    if exam_type == "True/False" or exam_type == "Mixed":
        base_prompt += f"""• TF (True/False):
  - Absolutely unambiguous statements
  - No "always/never" qualifiers unless truly accurate
  - {difficulty} example: "{diff_attr["examples"]["TF"]}"
  - Answer must be exactly "True" or "False"

"""

    if exam_type == "Identification" or exam_type == "Mixed":
        base_prompt += f"""• ID (Identification):
  - Questions with specific, precise term answers
  - Answers must be 1-5 words maximum
  - {difficulty} example: "{diff_attr["examples"]["ID"]}"
  - Answer should be a specific term or concept name

"""

    if exam_type == "Enumeration" or exam_type == "Mixed":
        base_prompt += f"""• ENUM (Enumeration):
  - Requests for specific lists of items or concepts
  - Answer as comma-separated list of specific terms
  - {difficulty} example: "{diff_attr["examples"]["ENUM"]}"
  - Be clear about the number of items expected

"""

    # Add additional difficulty-specific guidelines
    base_prompt += f"""
CRITICAL GUIDELINES FOR {difficulty.upper()} QUESTIONS:
"""

    if difficulty == "Easy":
        base_prompt += """• Focus on foundational concepts and basic terminology
• Use direct questions with clear-cut answers
• Test recognition and recall of essential information
• Keep answer options distinct and unambiguous (for MC)
• Avoid complex scenarios or applications
"""
    elif difficulty == "Medium":
        base_prompt += """• Test understanding of relationships between concepts
• Require application of knowledge to somewhat familiar situations
• Include questions that need moderate analysis
• Create plausible distractors based on partial understanding (for MC)
• Test ability to explain processes or compare related concepts
"""
    else:  # Hard
        base_prompt += """• Require synthesis of multiple concepts and deep understanding
• Test application in complex or novel scenarios
• Include challenging distractors that test subtle distinctions (for MC)
• Assess ability to evaluate, critique, or predict outcomes
• Require sophisticated reasoning and thorough content mastery
"""

    # Add content for generation
    base_prompt += f"""
Content to base questions on:
{module_text}

IMPORTANT: Generate exactly {num_q} different, well-formed questions of {difficulty} difficulty level following the above guidelines. Each question must have the correct structure based on its type.
"""

    return base_prompt


# ------------------------------ FLASHCARD GENERATION HELPERS ------------------------------

def clean_json_output(ai_text):
    """
    Attempts to extract and clean a JSON substring from the AI text.
    Handles both array responses and fixes common JSON formatting issues.
    """
    try:
        # Find the first '[' and last ']' to extract JSON array
        start_index = ai_text.find('[')
        end_index = ai_text.rfind(']') + 1
        
        if start_index >= 0 and end_index > start_index:
            json_str = ai_text[start_index:end_index]
            
            # Clean up common issues
            json_str = json_str.replace("\n", " ").replace("\r", " ")
            
            # Fix trailing commas (invalid in JSON)
            json_str = re.sub(r",\s*([\]}])", r"\1", json_str)
            
            # Fix missing quotes around property names
            json_str = re.sub(r'([{,])\s*([a-zA-Z_][a-zA-Z0-9_]*)\s*:', r'\1"\2":', json_str)
            
            # Fix single quotes used instead of double quotes
            in_string = False
            result = []
            i = 0
            while i < len(json_str):
                if json_str[i] == '"':
                    in_string = not in_string
                elif json_str[i] == "'" and not in_string:
                    result.append('"')
                    i += 1
                    continue
                result.append(json_str[i])
                i += 1
            json_str = ''.join(result)
            
            return json_str.strip()
        else:
            print("Error: Could not find JSON array markers in response")
            return "[]"  # Return empty array as fallback
    except Exception as e:
        print("Error cleaning JSON output:", e)
        return "[]"  # Return empty array as fallback


def generate_flashcards_from_text(text, topic):
    """
    Splits text into manageable chunks, calls the AI to generate flashcards for each chunk,
    cleans the JSON response, and returns a list of flashcards with unified keys:
      - question: The flashcard question
      - answer: The flashcard answer (specific terms)
    """
    chunk_size = 3000
    words = text.split()
    chunks = []
    current_chunk = []
    current_length = 0

    for word in words:
        if current_length + len(word) + 1 > chunk_size:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]
            current_length = len(word)
        else:
            current_chunk.append(word)
            current_length += len(word) + 1

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    flashcards = []
    for i, chunk in enumerate(chunks, start=1):
        prompt = f"""As an advanced educational AI specializing in precise, concise, term-focused flashcards, your task is to extract maximum learning value from content by creating exceptional term-only flashcards.

FLASHCARD JSON FORMAT:
{{
  "question": "<clear, precise, well-formed question>",
  "answer": "<ONLY specific terms or concise comma-separated lists>",
  "type": "ID or ENUM"
}}

SUPER-CRITICAL REQUIREMENTS (ABSOLUTE MUST-FOLLOW):
1. ANSWERS MUST CONTAIN ONLY TERMS, NEVER EXPLANATIONS OR SENTENCES
   - ID type: exact term (noun/name) limited to 1-5 words maximum
   - ENUM type: comma-separated list of specific terms only

2. PROHIBITED ANSWER PATTERNS (AUTOMATIC REJECTION):
   - Any explanation phrases like "refers to", "defined as", "a process", etc.
   - Any articles like "a", "an", "the" at the beginning
   - Any verbs like "is", "are", "means", "represents"
   - Any full sentence or phrase that's not a direct term

3. MANDATORY MIX OF TYPES (40% ENUM, 60% ID):
   - IF the text contains lists, classifications, categories = ENUM
   - IF the text defines key terms, concepts, specific names = ID
   - Create at least 3-4 ENUM questions regardless of content

4. MAXIMIZING COVERAGE OF ALL KEY TERMS:
   - Extract ALL important concepts, not just obvious ones
   - Cover key definitions, processes, classifications, examples
   - Ensure questions span entire content, not just portions
   - Include technical terms, proper names, and specialized vocabulary

EXAMPLES OF PERFECT ANSWERS:

ID ANSWERS (term only, no explanations):
✓ "Binary Search Tree"
✓ "Encapsulation"
✓ "HTTP"
✓ "Python"

ENUM ANSWERS (comma-separated terms only):
✓ "HTML, CSS, JavaScript"
✓ "Inheritance, Polymorphism, Encapsulation, Abstraction"
✓ "Stack, Queue, Linked List, Tree, Graph"

EXAMPLES OF REJECTED ANSWERS:
✗ "A data structure that uses LIFO ordering"
✗ "The process of hiding implementation details"
✗ "Refers to the transfer protocol used in web"
✗ "A programming language that is interpreted"

PERFECT FLASHCARD EXAMPLES:

ID EXAMPLES:
{{
  "question": "What data structure stores elements in Last-In-First-Out order?",
  "answer": "Stack",
  "type": "ID"
}}

ENUM EXAMPLES:
{{
  "question": "What are the primary data types in JavaScript?",
  "answer": "String, Number, Boolean, Object, Undefined, Null, Symbol, BigInt",
  "type": "ENUM"
}}

Generate 8-15 flashcards from the content below. REMEMBER: ANSWERS MUST BE TERMS ONLY (nouns/names), NEVER EXPLANATIONS. Focus on maximum extraction of useful knowledge.

Topic: {topic}
Text (chunk {i}/{len(chunks)}): {chunk}"""
        try:
            model = genai.GenerativeModel(MODEL_NAME)
            response = model.generate_content(
                prompt,
                generation_config=genai.types.GenerationConfig(
                    candidate_count=1, 
                    max_output_tokens=3072,
                    temperature=0.2,
                    top_p=0.95
                )
            )
            if not response or not hasattr(response, "text"):
                continue

            ai_text = response.text
            cleaned_text = clean_json_output(ai_text)
            fc = json.loads(cleaned_text)

            if isinstance(fc, list):
                for card in fc:
                    # Validate and clean the flashcard format
                    if "question" in card and "answer" in card:
                        question = card["question"].strip()
                        answer = card["answer"].strip()
                        card_type = card.get("type", "ID").strip().upper()
                        
                        # Enhanced rigorous answer validation
                        # 1. Check for obvious explanatory phrases and patterns
                        explanation_patterns = [
                            " is ", " are ", " refers ", " defined ", " process ", " means ", 
                            " concept ", " type ", " example ", " approach ", " method ",
                            " a ", " an ", " the ", " where ", " when ", " how ", " why ",
                            " used ", " allows ", " enables ", " represents ", " involves ",
                            " consists ", " contains ", " includes ", " describes ", " technique ",
                            " that ", " which ", " who ", " whom ", " whose "
                        ]
                        
                        if any(pattern in " " + answer.lower() + " " for pattern in explanation_patterns):
                            continue
                        
                        # 2. Type-specific stricter validation
                        if card_type == "ID":
                            # ID answer must be 1-5 words, no commas
                            if len(answer.split()) > 5 or "," in answer:
                                continue
                            # Reject if complete sentence structure detected (has verb)
                            verbs = ["is", "are", "was", "were", "be", "being", "been", "have", "has", "had",
                                    "do", "does", "did", "can", "could", "will", "would", "shall", "should",
                                    "may", "might", "must", "refers", "means", "represents", "consists",
                                    "involves", "allows", "enables", "helps", "works", "functions"]
                            if any(f" {verb} " in f" {answer.lower()} " for verb in verbs):
                                continue
                        elif card_type == "ENUM":
                            # ENUM must have commas separating terms
                            if "," not in answer:
                                continue
                            # Each term must be a specific term (not explanation/sentence)
                            terms = [t.strip() for t in answer.split(",")]
                            if any(len(t.split()) > 5 for t in terms):
                                continue
                            # Reject if any term contains explanatory language
                            if any(any(pattern in " " + term.lower() + " " for pattern in explanation_patterns) for term in terms):
                                continue
                        else:
                            # Unknown card type, apply generic validation
                            if len(answer.split()) > 5 and "," not in answer:
                                continue
                        
                        # 3. Ensure question quality
                        # Skip questions that are too short or don't end with a question mark
                        if len(question) < 15 or not question.endswith("?"):
                            continue
                            
                        # Skip empty questions/answers or pure numbers as answers
                        if question and answer and not answer.isdigit() and len(question) > 10:
                            flashcards.append({
                                "question": question,
                                "answer": answer
                            })
        except Exception as e:
            print("Error generating flashcards:", e)
    
    # Ensure we have a good mix of card types and avoid duplicates
    seen_answers = set()
    unique_flashcards = []
    for card in flashcards:
        answer_key = card["answer"].lower().strip()
        # Don't add exact duplicates
        if answer_key not in seen_answers:
            seen_answers.add(answer_key)
            unique_flashcards.append(card)
    
    # Analyze and ensure we have both ID and ENUM type cards
    enum_cards = [card for card in unique_flashcards if "," in card["answer"]]
    id_cards = [card for card in unique_flashcards if "," not in card["answer"]]
    
    # If we're missing enumeration cards, try one more pass to generate some
    if len(enum_cards) < min(3, len(unique_flashcards) * 0.25):
        for i, chunk in enumerate(chunks, start=1):
            try:
                enum_prompt = f"""Generate ONLY enumeration flashcards with comma-separated lists of terms as answers from this content. Each answer must be ONLY a list of specific terms (2-8 terms), not explanations.

Example good flashcard:
{{
  "question": "What are the four main data types in SQL?",
  "answer": "INTEGER, VARCHAR, DATE, BOOLEAN"
}}

Content for enumeration flashcards:
{chunk}"""
                
                model = genai.GenerativeModel(MODEL_NAME)
                response = model.generate_content(
                    enum_prompt,
                    generation_config=genai.types.GenerationConfig(candidate_count=1, temperature=0.3)
                )
                if response and hasattr(response, "text"):
                    enum_text = response.text
                    cleaned_enum = clean_json_output(enum_text)
                    try:
                        enum_fc = json.loads(cleaned_enum)
                        if isinstance(enum_fc, list):
                            for card in enum_fc:
                                if "question" in card and "answer" in card:
                                    q = card["question"].strip()
                                    a = card["answer"].strip()
                                    if q and a and "," in a and q.endswith("?") and all(len(term.strip().split()) <= 5 for term in a.split(",")):
                                        unique_flashcards.append({"question": q, "answer": a})
                    except:
                        pass
            except:
                pass
            
            # If we've got enough enum cards now, stop
            if len([c for c in unique_flashcards if "," in c["answer"]]) >= 3:
                break
    
    return unique_flashcards


# --------------------------- ROUTES ------------------------------

@app.route('/')
def index():
    # Bypass login and directly access dashboard
    return redirect(url_for('dashboard'))


@app.route('/dashboard')
def dashboard():
    return render_template('dashboard.html')


# SUBJECT MANAGEMENT ROUTES
@app.route('/subjects')
def subjects():
    all_subjects = Subject.query.all()
    return render_template('subjects.html', subjects=all_subjects)


@app.route('/add_subject', methods=['POST'])
def add_subject():
    subject_name = request.form.get('subjectName')
    if not subject_name:
        flash('Subject name cannot be empty!', 'error')
        return redirect(url_for('subjects'))
    
    existing_subject = Subject.query.filter_by(name=subject_name).first()
    if existing_subject:
        flash(f'Subject "{subject_name}" already exists!', 'error')
        return redirect(url_for('subjects'))
    
    new_subject = Subject(name=subject_name)
    db.session.add(new_subject)
    db.session.commit()
    flash(f'Subject "{subject_name}" added successfully!', 'success')
    return redirect(url_for('subjects'))


@app.route('/delete_subject/<int:subject_id>', methods=['POST'])
def delete_subject(subject_id):
    subject = Subject.query.get_or_404(subject_id)
    subject_name = subject.name
    db.session.delete(subject)
    db.session.commit()
    flash(f'Subject "{subject_name}" and all its topics have been deleted.', 'success')
    return redirect(url_for('subjects'))


@app.route('/subject/<int:subject_id>')
def view_subject(subject_id):
    subject = Subject.query.get_or_404(subject_id)
    topics = Topic.query.filter_by(subject_id=subject_id).all()
    quiz_results = QuizResult.query.filter_by(subject_id=subject_id).order_by(QuizResult.created_at.desc()).limit(10).all()
    return render_template('topics.html', subject=subject, topics=topics, quiz_results=quiz_results)


@app.route('/add_topic/<int:subject_id>', methods=['POST'])
def add_topic(subject_id):
    subject = Subject.query.get_or_404(subject_id)
    topic_name = request.form.get('topicName')
    
    if not topic_name:
        flash('Topic name cannot be empty!', 'error')
        return redirect(url_for('view_subject', subject_id=subject_id))
    
    existing_topic = Topic.query.filter_by(name=topic_name, subject_id=subject_id).first()
    if existing_topic:
        flash(f'Topic "{topic_name}" already exists in this subject!', 'error')
        return redirect(url_for('view_subject', subject_id=subject_id))
    
    new_topic = Topic(name=topic_name, subject_id=subject_id)
    db.session.add(new_topic)
    db.session.commit()
    flash(f'Topic "{topic_name}" added successfully!', 'success')
    return redirect(url_for('view_subject', subject_id=subject_id))


@app.route('/delete_topic/<int:topic_id>', methods=['POST'])
def delete_topic(topic_id):
    topic = Topic.query.get_or_404(topic_id)
    subject_id = topic.subject_id
    topic_name = topic.name
    db.session.delete(topic)
    db.session.commit()
    flash(f'Topic "{topic_name}" and all its flashcards have been deleted.', 'success')
    return redirect(url_for('view_subject', subject_id=subject_id))


@app.route('/view_flashcards/<int:topic_id>')
def view_flashcards(topic_id):
    topic = Topic.query.get_or_404(topic_id)
    flashcards = Flashcard.query.filter_by(topic_id=topic_id).all()
    
    if not flashcards:
        flash('No flashcards found for this topic.', 'info')
        return redirect(url_for('view_subject', subject_id=topic.subject_id))
    
    # Convert database flashcards to format expected by the template
    formatted_flashcards = [{'question': card.question, 'answer': card.answer} for card in flashcards]
    
    return render_template('flashcard_view.html', 
                          flashcards=formatted_flashcards, 
                          subject=topic.subject.name, 
                          topic=topic.name,
                          topic_id=topic_id)


@app.route('/help')
def help_page():
    return render_template('help.html')


@app.route('/about')
def about_page():
    return render_template('about.html')


@app.route('/contact')
def contact_page():
    return render_template('contact_us.html')


@app.route('/static/<path:filename>')
def static_files(filename):
    return send_file(os.path.join('static', filename))


@app.route('/texttospeech', methods=['GET', 'POST'])
def text_to_speech():
    return render_template('text_to_speech.html')


@app.route('/timer')
def timer():
    return render_template('timer.html')


@app.route('/quiz_options')
@app.route('/quiz_options/<int:subject_id>')
def quiz_options(subject_id=None):
    if subject_id is None:
        # Old route without subject specified, redirect to subjects page
        return redirect(url_for('subjects'))
    
    subject = Subject.query.get_or_404(subject_id)
    topics = Topic.query.filter_by(subject_id=subject_id).all()
    return render_template('quiz_options.html', subject=subject, topics=topics)


@app.route('/quizmaker_ai', methods=['GET', 'POST'])
def quiz():
    if request.method == 'POST':
        subject_id = request.form.get('subject_id')
        topic_id = request.form.get('topic_id')
        
        if not subject_id:
            flash('Please select a subject first!', 'error')
            return redirect(url_for('subjects'))
            
        subject = Subject.query.get_or_404(int(subject_id))
        topic = None
        if topic_id:
            topic = Topic.query.get_or_404(int(topic_id))
            
        question_type = request.form.get('questionType', 'Multiple Choice')
        difficulty = request.form.get('difficulty', 'Medium')
        num_questions = int(request.form.get('numQuestions', 5))
        
        # Process uploaded file
        uploaded_file = request.files.get('fileUpload')
        module_text = ""
        if uploaded_file and uploaded_file.filename:
            ext = os.path.splitext(uploaded_file.filename)[1].lower()
            if ext == ".pdf":
                try:
                    reader = PdfReader(uploaded_file.stream)
                    text_chunks = []
                    for page in reader.pages:
                        page_text = page.extract_text() or ""
                        text_chunks.append(page_text)
                    module_text = "\n".join(text_chunks)
                except Exception as e:
                    flash("Error reading PDF file: " + str(e), "error")
                    return redirect(url_for('quiz_options'))
            elif ext == ".docx":
                try:
                    doc = docx.Document(uploaded_file)
                    module_text = "\n".join([para.text for para in doc.paragraphs])
                except Exception as e:
                    flash("Error reading DOCX file: " + str(e), "error")
                    return redirect(url_for('quiz_options'))
            elif ext in [".ppt", ".pptx"]:
                try:
                    prs = Presentation(uploaded_file)
                    text_chunks = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_chunks.append(shape.text)
                    module_text = "\n".join(text_chunks)
                except Exception as e:
                    flash("Error reading PPT file: " + str(e), "error")
                    return redirect(url_for('quiz_options'))
            elif ext == ".txt":
                try:
                    file_data = uploaded_file.read()
                    module_text = file_data.decode("utf-8", errors="ignore")
                except Exception as e:
                    flash("Error reading TXT file: " + str(e), "error")
                    return redirect(url_for('quiz_options'))
            else:
                flash("Unsupported file format", "error")
                return redirect(url_for('quiz_options'))
        if not module_text.strip():
            flash("No module content found. Please provide valid module content.", "error")
            return redirect(url_for('quiz_options'))

        # Set temperature based on difficulty for appropriate question complexity
        temperature = 0.5  # base value
        if difficulty == "Easy":
            temperature = 0.3  # more consistent, straightforward questions
        elif difficulty == "Medium":
            temperature = 0.5  # balanced approach
        elif difficulty == "Hard":
            temperature = 0.7  # more creative, challenging questions

        # Build prompt and generate quiz using AI
        prompt = build_quiz_prompt(module_text, question_type, difficulty, num_questions)
        try:
            model = genai.GenerativeModel(MODEL_NAME)
            response = model.generate_content(
                prompt,
                generation_config=genai.types.GenerationConfig(
                    candidate_count=1,
                    max_output_tokens=4096,
                    temperature=temperature,
                    top_p=0.95,
                    top_k=40,
                    response_mime_type="application/json"
                )
            )
            if not response or not hasattr(response, "text"):
                raise ValueError("Invalid AI response.")
                
            ai_text = response.text.strip()
            
            # Clean the response and ensure it's valid JSON
            cleaned_json = clean_json_output(ai_text)
            
            # Parse the AI content into quiz questions
            quiz_data = parse_ai_content(cleaned_json, num_questions)
            if not quiz_data['questions']:
                # If we didn't get any questions, try a second attempt with more structure
                retry_prompt = f"Generate {num_questions} {question_type} questions about {module_text[:300]}... Please follow a strict JSON format with question, answer, type, and options fields. Make them {difficulty.lower()} difficulty."
                retry_response = model.generate_content(
                    retry_prompt,
                    generation_config=genai.types.GenerationConfig(
                        candidate_count=1,
                        max_output_tokens=4096,
                        temperature=0.7,
                        response_mime_type="application/json"
                    )
                )
                
                if retry_response and hasattr(retry_response, "text"):
                    retry_json = clean_json_output(retry_response.text.strip())
                    quiz_data = parse_ai_content(retry_json, num_questions)
                
                if not quiz_data['questions']:
                    flash("AI could not generate valid questions. Please try again with different content or settings.", "error")
                    return redirect(url_for('quiz_options'))
        except Exception as e:
            flash(f"AI generation failed: {str(e)}. Please try again.", "error")
            return redirect(url_for('quiz_options'))

        # Store information about the quiz in session
        session['quiz_subject_id'] = subject_id
        session['quiz_topic_id'] = topic_id
        session['quiz_question_type'] = question_type
        session['quiz_difficulty'] = difficulty
        
        return render_template('quiz.html', 
                              questions=quiz_data['questions'], 
                              answers=quiz_data['answers'],
                              options=quiz_data['options'], 
                              subject=subject.name,
                              topic=topic.name if topic else None, 
                              question_types=quiz_data['question_types'])
    return redirect(url_for('quiz_options'))


@app.route('/submit_quiz', methods=['POST'])
def submit_quiz():
    try:
        data = request.get_json()
        score = data.get('score', 0)
        total = data.get('total', 0)
        
        # Retrieve quiz information from session
        subject_id = session.get('quiz_subject_id')
        topic_id = session.get('quiz_topic_id')
        question_type = session.get('quiz_question_type')
        difficulty = session.get('quiz_difficulty')
        
        if not subject_id:
            return jsonify({'status': 'error', 'message': 'Subject information missing'})
        
        # Save quiz result to database
        quiz_result = QuizResult(
            score=score,
            total_questions=total,
            question_type=question_type,
            difficulty=difficulty,
            subject_id=subject_id,
            topic_id=topic_id if topic_id else None
        )
        db.session.add(quiz_result)
        db.session.commit()
        
        return jsonify({'status': 'success', 'message': 'Quiz results saved successfully'})
    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)})


@app.route('/delete_quiz_result/<int:result_id>', methods=['POST'])
def delete_quiz_result(result_id):
    quiz_result = QuizResult.query.get_or_404(result_id)
    subject_id = quiz_result.subject_id
    
    try:
        db.session.delete(quiz_result)
        db.session.commit()
        flash('Quiz result deleted successfully.', 'success')
    except Exception as e:
        flash(f'Error deleting quiz result: {str(e)}', 'error')
    
    return redirect(url_for('view_subject', subject_id=subject_id))


@app.route('/flashcard', methods=['GET', 'POST'])
@app.route('/flashcard/<int:topic_id>', methods=['GET', 'POST'])
def flashcard(topic_id=None):
    if topic_id is None:
        # Old route without topic specified, redirect to subjects page
        return redirect(url_for('subjects'))
        
    topic = Topic.query.get_or_404(topic_id)
    
    if request.method == 'POST':
        # Get file upload or text input
        file = request.files.get('fileUpload')
        text_input = request.form.get('textInput')
        
        if not file and not text_input:
            flash('Please provide either a file or text input.', 'error')
            return redirect(url_for('flashcard', topic_id=topic_id))
            
        # Extract text from upload or use text input directly
        if file:
            ext = os.path.splitext(file.filename)[1].lower()
            if ext in ['.pdf', '.docx', '.ppt', '.pptx']:
                try:
                    if ext == '.pdf':
                        reader = PdfReader(file.stream)
                        text_chunks = []
                        for page in reader.pages:
                            page_text = page.extract_text() or ""
                            text_chunks.append(page_text)
                        text_input = "\n".join(text_chunks)
                    elif ext == '.docx':
                        doc = docx.Document(file)
                        text_input = "\n".join([para.text for para in doc.paragraphs])
                    elif ext in ['.ppt', '.pptx']:
                        prs = Presentation(file)
                        text_chunks = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text_chunks.append(shape.text)
                        text_input = "\n".join(text_chunks)
                except Exception as e:
                    flash("Error reading file: " + str(e), "error")
                    return redirect(url_for('flashcard', topic_id=topic_id))
            else:
                flash("Unsupported file format", "error")
                return redirect(url_for('flashcard', topic_id=topic_id))

        if not text_input.strip():
            flash("Please provide input text or upload a file.", "error")
            return redirect(url_for('flashcard', topic_id=topic_id))

        try:
            flashcards = generate_flashcards_from_text(text_input, topic.name)
            if not flashcards:
                flash("No flashcards were generated. Try again with more specific content.", "error")
                return redirect(url_for('flashcard', topic_id=topic_id))

            # Save flashcards to database
            for card in flashcards:
                new_flashcard = Flashcard(
                    question=card['question'],
                    answer=card['answer'],
                    topic_id=topic_id
                )
                db.session.add(new_flashcard)
            
            db.session.commit()
            flash('Flashcards created successfully!', 'success')
            
            # Redirect to view the created flashcards
            return redirect(url_for('view_flashcards', topic_id=topic_id))
        except Exception as e:
            flash(f"Error generating flashcards: {str(e)}", "error")
            return redirect(url_for('flashcard', topic_id=topic_id))
    else:
        return render_template('flashcard_upload.html', topic=topic)


@app.route('/save_flashcards', methods=['POST'])
def save_flashcards():
    flashcards = session.get('flashcards', [])
    subject = session.get('subject', 'default_subject')
    topic = session.get('topic', 'default_topic')
    if not flashcards:
        flash("No flashcards to save.", "error")
        return redirect(url_for('flashcard'))

    def sanitize_string(s):
        return re.sub(r'\s+', '_', re.sub(r'[^a-z0-9\s]', '', s.lower())).strip()

    subject_clean = sanitize_string(subject)
    topic_clean = sanitize_string(topic)
    file_name = f"{subject_clean}_{topic_clean}_flashcards.txt"

    lines = []
    for idx, card in enumerate(flashcards, start=1):
        lines.append(f"Flashcard {idx}")
        lines.append(f"Q: {card['question']}")
        lines.append(f"A: {card['answer']}")
        lines.append("")
    text_data = "\n".join(lines)

    memory_file = BytesIO()
    memory_file.write(text_data.encode('utf-8'))
    memory_file.seek(0)

    return send_file(
        memory_file,
        as_attachment=True,
        download_name=file_name,
        mimetype="text/plain"
    )


@app.route('/logout')
def logout():
    session.clear()
    flash("You have been logged out.", "info")
    return redirect(url_for('login'))


@app.errorhandler(404)
def not_found_error(error):
    try:
        return render_template('404.html'), 404
    except TemplateNotFound:
        return "<h1>404 Not Found</h1>", 404


@app.errorhandler(500)
def internal_error(error):
    try:
        return render_template('500.html'), 500
    except TemplateNotFound:
        return "<h1>500 Internal Server Error</h1>", 500


if __name__ == '__main__':
    app.run(debug=True)
